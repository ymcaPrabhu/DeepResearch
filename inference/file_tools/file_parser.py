import json
import os
import re
import time
import zipfile
import math
from pathlib import Path

from typing import Any, Dict, List, Optional, Union
from collections import Counter
import xml.etree.ElementTree as ET
from pandas import Timestamp
from datetime import datetime
from pandas.api.types import is_datetime64_any_dtype

import pandas as pd
from tabulate import tabulate
from qwen_agent.log import logger
from qwen_agent.settings import DEFAULT_WORKSPACE, DEFAULT_MAX_INPUT_TOKENS
from qwen_agent.tools.base import BaseTool, register_tool
from qwen_agent.tools.storage import KeyNotExistsError, Storage
from file_tools.utils import (get_file_type, hash_sha256, is_http_url, get_basename_from_url, 
                                  sanitize_chrome_file_path, save_url_to_local_work_dir)
from qwen_agent.utils.tokenization_qwen import count_tokens, tokenizer
from file_tools.idp import IDP

# Configuration constants
PARSER_SUPPORTED_FILE_TYPES = ['pdf', 'docx', 'pptx', 'txt', 'html', 'csv', 'tsv', 'xlsx', 'xls', 'doc', 'zip', '.mp4', '.mov', '.mkv', '.webm', '.mp3', '.wav']
USE_IDP = os.getenv("USE_IDP", True)
IDP_TIMEOUT = 150000
ENABLE_CSI = False
PARAGRAPH_SPLIT_SYMBOL = '\n'


class CustomJSONEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, (datetime, Timestamp)):
            return obj.isoformat()
        return super().default(obj)


class FileParserError(Exception):
    """Custom exception for document parsing errors"""

    def __init__(self, message: str, code: str = '400', exception: Optional[Exception] = None):
        super().__init__(message)
        self.code = code
        self.exception = exception


def parse_file_by_idp(file_path: str = None, file_url: str = None) -> List[dict]:
    idp = IDP()
    try:
        fid = idp.file_submit_with_url(file_url) if file_url else idp.file_submit_with_path(file_path)
        if not fid:
            return []

        for _ in range(10):  
            result, status = idp.file_parser_query(fid)
            if status == 'success':
                return process_idp_result(result)
            time.sleep(10)

        logger.error("IDP parsing timeout")
        return []
    except Exception as e:
        logger.error(f"IDP processing failed: {str(e)}")
        return []


def process_idp_result(result: dict) -> List[dict]:
    pages = []
    current_page = None

    for layout in result.get('layouts', []):
        page_num = layout.get('pageNum', 0)
        content = layout.get('markdownContent', '')

        if current_page and current_page['page_num'] == page_num:
            current_page['content'].append({'text': content})
        else:
            current_page = {'page_num': page_num, 'content': [{'text': content}]}
            pages.append(current_page)

    return pages


def clean_text(text: str) -> str:
    cleaners = [
        lambda x: re.sub(r'\n+', '\n', x),  
        lambda x: x.replace("Add to Qwen's Reading List", ''),
        lambda x: re.sub(r'-{6,}', '-----', x),  
        lambda x: x.strip()
    ]
    for cleaner in cleaners:
        text = cleaner(text)
    return text


def get_plain_doc(doc: list):
    paras = []
    for page in doc:
        for para in page['content']:
            for k, v in para.items():
                if k in ['text', 'table', 'image']:
                    paras.append(v)
    return PARAGRAPH_SPLIT_SYMBOL.join(paras)


def df_to_markdown(df: pd.DataFrame) -> str:
    df = df.dropna(how='all').fillna('')
    return tabulate(df, headers='keys', tablefmt='pipe', showindex=False)


def parse_word(docx_path: str, extract_image: bool = False):
    if extract_image:
        raise ValueError('Currently, extracting images is not supported!')

    from docx import Document
    doc = Document(docx_path)

    content = []
    for para in doc.paragraphs:
        content.append({'text': para.text})
    for table in doc.tables:
        tbl = []
        for row in table.rows:
            tbl.append('|' + '|'.join([cell.text for cell in row.cells]) + '|')
        tbl = '\n'.join(tbl)
        content.append({'table': tbl})
    return [{'page_num': 1, 'content': content}]


def parse_ppt(path: str, extract_image: bool = False):
    if extract_image:
        raise ValueError('Currently, extracting images is not supported!')

    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
    try:
        ppt = Presentation(path)
    except PackageNotFoundError as ex:
        logger.warning(ex)
        return []
    doc = []
    for slide_number, slide in enumerate(ppt.slides):
        page = {'page_num': slide_number + 1, 'content': []}

        for shape in slide.shapes:
            if not shape.has_text_frame and not shape.has_table:
                pass

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    paragraph_text = clean_text(paragraph_text)
                    if paragraph_text.strip():
                        page['content'].append({'text': paragraph_text})

            if shape.has_table:
                tbl = []
                for row_number, row in enumerate(shape.table.rows):
                    tbl.append('|' + '|'.join([cell.text for cell in row.cells]) + '|')
                tbl = '\n'.join(tbl)
                page['content'].append({'table': tbl})
        doc.append(page)
    return doc


def parse_pdf(pdf_path: str, extract_image: bool = False) -> List[dict]:
    # Todo: header and footer
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTImage, LTRect, LTTextContainer

    doc = []
    import pdfplumber
    pdf = pdfplumber.open(pdf_path)
    for i, page_layout in enumerate(extract_pages(pdf_path)):
        page = {'page_num': page_layout.pageid, 'content': []}

        elements = []
        for element in page_layout:
            elements.append(element)

        # Init params for table
        table_num = 0
        tables = []

        for element in elements:
            if isinstance(element, LTRect):
                if not tables:
                    tables = extract_tables(pdf, i)
                if table_num < len(tables):
                    table_string = table_converter(tables[table_num])
                    table_num += 1
                    if table_string:
                        page['content'].append({'table': table_string, 'obj': element})
            elif isinstance(element, LTTextContainer):
                # Delete line breaks in the same paragraph
                text = element.get_text()
                # Todo: Further analysis using font
                font = get_font(element)
                if text.strip():
                    new_content_item = {'text': text, 'obj': element}
                    if font:
                        new_content_item['font-size'] = round(font[1])
                        # new_content_item['font-name'] = font[0]
                    page['content'].append(new_content_item)
            elif extract_image and isinstance(element, LTImage):
                # Todo: ocr
                raise ValueError('Currently, extracting images is not supported!')
            else:
                pass

        # merge elements
        page['content'] = postprocess_page_content(page['content'])
        doc.append(page)

    return doc


def parse_txt(path: str):
    with open(path, 'r', encoding='utf-8') as f:  
        text = f.read()
    paras = text.split(PARAGRAPH_SPLIT_SYMBOL)
    content = []
    for p in paras:
        content.append({'text': p})
    return [{'page_num': 1, 'content': content}]


def get_font(element):
    from pdfminer.layout import LTChar, LTTextContainer

    fonts_list = []
    for text_line in element:
        if isinstance(text_line, LTTextContainer):
            for character in text_line:
                if isinstance(character, LTChar):
                    fonts_list.append((character.fontname, character.size))

    fonts_list = list(set(fonts_list))
    if fonts_list:
        counter = Counter(fonts_list)
        most_common_fonts = counter.most_common(1)[0][0]
        return most_common_fonts
    else:
        return []


def extract_tables(pdf, page_num):
    table_page = pdf.pages[page_num]
    tables = table_page.extract_tables()
    return tables


def table_converter(table):
    table_string = ''
    for row_num in range(len(table)):
        row = table[row_num]
        cleaned_row = [
            item.replace('\n', ' ') if item is not None and '\n' in item else 'None' if item is None else item
            for item in row
        ]
        table_string += ('|' + '|'.join(cleaned_row) + '|' + '\n')
    table_string = table_string[:-1]
    return table_string


def postprocess_page_content(page_content: list) -> list:
    # rm repetitive identification for table and text
    # Some documents may repeatedly recognize LTRect and LTTextContainer
    table_obj = [p['obj'] for p in page_content if 'table' in p]
    tmp = []
    for p in page_content:
        repetitive = False
        if 'text' in p:
            for t in table_obj:
                if t.bbox[0] <= p['obj'].bbox[0] and p['obj'].bbox[1] <= t.bbox[1] and t.bbox[2] <= p['obj'].bbox[
                    2] and p['obj'].bbox[3] <= t.bbox[3]:
                    repetitive = True
                    break

        if not repetitive:
            tmp.append(p)
    page_content = tmp

    # merge paragraphs that have been separated by mistake
    new_page_content = []
    for p in page_content:
        if new_page_content and 'text' in new_page_content[-1] and 'text' in p and abs(
                p.get('font-size', 12) -
                new_page_content[-1].get('font-size', 12)) < 2 and p['obj'].height < p.get('font-size', 12) + 1:
            # Merge those lines belonging to a paragraph
            new_page_content[-1]['text'] += f' {p["text"]}'
            # new_page_content[-1]['font-name'] = p.get('font-name', '')
            new_page_content[-1]['font-size'] = p.get('font-size', 12)
        else:
            p.pop('obj')
            new_page_content.append(p)
    for i in range(len(new_page_content)):
        if 'text' in new_page_content[i]:
            new_page_content[i]['text'] = clean_text(new_page_content[i]['text'])
    return new_page_content


def extract_xls_schema(file_path: str) -> Dict[str, Any]:
    xls = pd.ExcelFile(file_path)
    schema = {
        "sheets": [],
        "n_sheets": len(xls.sheet_names)
    }

    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name, nrows=3)  # 读取前3行

        dtype_mapping = {
            'object': 'string',
            'datetime64[ns]': 'datetime',
            'timedelta64[ns]': 'timedelta'
        }
        dtypes = df.dtypes.astype(str).replace(dtype_mapping).to_dict()

        sample_df = df.head(3).copy()
        for col in sample_df.columns:
            if is_datetime64_any_dtype(sample_df[col]):
                sample_df[col] = sample_df[col].dt.strftime('%Y-%m-%dT%H:%M:%S')

        sheet_info = {
            "name": sheet_name,
            "columns": df.columns.tolist(),
            "dtypes": dtypes,  
            "sample_data": sample_df.to_dict(orient='list') 
        }
        schema["sheets"].append(sheet_info)

    return schema


def extract_csv_schema(file_path: str) -> Dict[str, Any]:
    df_dtype = pd.read_csv(file_path, nrows=100)  
    df_sample = pd.read_csv(file_path, nrows=3) 

    return {
        "columns": df_dtype.columns.tolist(),
        "dtypes": df_dtype.dtypes.astype(str).to_dict(),
        "sample_data": df_sample.to_dict(orient='list'),
        "estimated_total_rows": _estimate_total_rows(file_path)
    }


def _estimate_total_rows(file_path) -> int:
    with open(file_path, 'rb') as f:
        line_count = 0
        chunk_size = 1024 * 1024  
        while chunk := f.read(chunk_size):
            line_count += chunk.count(b'\n')
    return line_count - 1  


def parse_tabular_file(file_path: str, **kwargs) -> List[dict]:
    try:
        df = pd.read_excel(file_path) if file_path.endswith(('.xlsx', '.xls')) else \
            pd.read_csv(file_path, **kwargs)
        if count_tokens(df_to_markdown(df)) > DEFAULT_MAX_INPUT_TOKENS:
            schema = extract_xls_schema(file_path) if file_path.endswith(('.xlsx', '.xls')) else \
                extract_csv_schema(file_path)
            return [{'page_num': 1, 'content': [{'schema': schema}]}]
        else:
            return [{'page_num': 1, 'content': [{'table': df_to_markdown(df)}]}]
    except Exception as e:
        logger.error(f"Table parsing failed: {str(e)}")
        return []  


def parse_zip(file_path: str, extract_dir: str) -> List[dict]:
    with zipfile.ZipFile(file_path, 'r') as zip_ref:
        zip_ref.extractall(extract_dir)
        return [os.path.join(extract_dir, f) for f in zip_ref.namelist()]


def parse_html(file_path: str) -> List[dict]:
    from bs4 import BeautifulSoup  

    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'lxml')

    content = [{'text': clean_text(p.get_text())}
               for p in soup.find_all(['p', 'div']) if p.get_text().strip()]

    return [{
        'page_num': 1,
        'content': content,
        'title': soup.title.string if soup.title else ''
    }]


def extract_xml_skeleton_markdown(xml_file):
    tree = ET.parse(xml_file)
    root = tree.getroot()
    markdown_lines = []

    def process_element(element, level=0, parent_path="", is_last=True, prefix=""):
        if level > 0:
            connector = "└── " if is_last else "├── "
            markdown_lines.append(f"{prefix}{connector}**{element.tag}**")
        else:
            markdown_lines.append(f"## Root: {element.tag}")

        if element.attrib:
            attrs = [f"`{k}`" for k in element.attrib.keys()]
            attr_line = f"{prefix}{'    ' if level > 0 else ''}*Attributes:* {', '.join(attrs)}"
            markdown_lines.append(attr_line)

        if element.text and element.text.strip():
            text_line = f"{prefix}{'    ' if level > 0 else ''}*Has text content*"
            markdown_lines.append(text_line)
        seen_tags = set()
        unique_children = []
        for child in element:
            if child.tag not in seen_tags:
                seen_tags.add(child.tag)
                unique_children.append(child)

        for i, child in enumerate(unique_children):
            is_last_child = (i == len(unique_children) - 1)
            child_prefix = prefix + ("    " if is_last else "│   ")
            process_element(child, level + 1,
                            f"{parent_path}/{element.tag}" if parent_path else element.tag,
                            is_last_child, child_prefix)

    process_element(root)
    markdown_content = "\n".join(markdown_lines)
    return markdown_content


def parse_xml(file_path: str) -> List[dict]:
    with open(file_path, 'r', encoding='utf-8') as f: 
        text = f.read()
    if count_tokens(text) > DEFAULT_MAX_INPUT_TOKENS:
        schema = extract_xml_skeleton_markdown(file_path)
        content = [{'schema': schema}]
    else:
        content = [{'text': text}]
    return [{'page_num': 1, 'content': content}]


def compress(results: list) -> list[str]:
    compress_results = []
    max_token = math.floor(DEFAULT_MAX_INPUT_TOKENS / len(results))
    for result in results:
        token_list = tokenizer.tokenize(result)
        token_list = token_list[:min(len(token_list), max_token)]
        compress_results.append(tokenizer.convert_tokens_to_string(token_list))
    return compress_results


# @register_tool('file_parser')
class SingleFileParser(BaseTool):
    name="file_parser"
    description = f"File parsing tool, supports parsing data in  {'/'.join(PARSER_SUPPORTED_FILE_TYPES)} formats, and returns the parsed markdown format data."
    parameters = [{
        'name': 'url',
        'type': 'string',
        'description': 'The full path of the file to be parsed, which can be a local path or a downloadable http(s) link.',
        'required': True
    }]

    def __init__(self, cfg: Optional[Dict] = None):
        super().__init__(cfg)
        self.data_root = self.cfg.get('path', os.path.join(DEFAULT_WORKSPACE, 'tools', self.name))
        self.db = Storage({'storage_root_path': self.data_root})
        self.structured_doc = self.cfg.get('structured_doc', True)

  
        self.parsers = {
            'pdf': parse_pdf,
            'docx': parse_word,
            'doc': parse_word,
            'pptx': parse_ppt,
            'txt': parse_txt,
            'jsonl': parse_txt,
            'jsonld': parse_txt,
            'pdb': parse_txt,
            'py': parse_txt,
            'html': parse_html,
            'xml': parse_xml,
            'csv': lambda p: parse_tabular_file(p, sep=','),
            'tsv': lambda p: parse_tabular_file(p, sep='\t'),
            'xlsx': parse_tabular_file,
            'xls': parse_tabular_file,
            'zip': self.parse_zip
        }

    def call(self, params: Union[str, dict], **kwargs) -> Union[str, list]:
        params = self._verify_json_format_args(params)
        file_path = self._prepare_file(params['url'])
        try:
            cached = self.db.get(f'{hash_sha256(file_path)}_ori')
            return self._flatten_result(json.loads(cached))
        except KeyNotExistsError:
            return self._flatten_result(self._process_new_file(file_path))

    def _prepare_file(self, path: str) -> str:
        if is_http_url(path):
            download_dir = os.path.join(self.data_root, hash_sha256(path))
            os.makedirs(download_dir, exist_ok=True)
            return save_url_to_local_work_dir(path, download_dir)
        return sanitize_chrome_file_path(path)

    def _process_new_file(self, file_path: str) -> Union[str, list]:
        file_type = get_file_type(file_path)
        idp_types = ['pdf', 'docx', 'pptx', 'xlsx', 'jpg', 'png', 'mp3']
        logger.info(f'Start parsing {file_path}...')
        logger.info(f'File type {file_type}...')
        logger.info(f"structured_doc {self.cfg.get('structured_doc')}...")

        if file_type not in idp_types:
            file_type = get_basename_from_url(file_path).split('.')[-1].lower()

        try:
            if USE_IDP and file_type in idp_types:
                try:
                    results = parse_file_by_idp(file_path=file_path)
                except Exception as e:
                    results = self.parsers[file_type](file_path)
            else:
                results = self.parsers[file_type](file_path)
            tokens = 0
            for page in results:
                for para in page['content']:
                    if 'schema' in para:
                        para['token'] = count_tokens(json.dumps(para['schema']))
                    else:
                        para['token'] = count_tokens(para.get('text', para.get('table')))
                    tokens += para['token']

            if not results or not tokens:
                logger.error(f"Parsing failed: No information was parsed")
                raise FileParserError("Document parsing failed")
            else:
                self._cache_result(file_path, results)
                return results
        except Exception as e:
            logger.error(f"Parsing failed: {str(e)}")
            raise FileParserError("Document parsing failed", exception=e)

    def _cache_result(self, file_path: str, result: list):
        cache_key = f'{hash_sha256(file_path)}_ori'
        self.db.put(cache_key, json.dumps(result, ensure_ascii=False))
        logger.info(f'The parsing result of {file_path} has been cached')

    def _flatten_result(self, result: list) -> str:
        return PARAGRAPH_SPLIT_SYMBOL.join(
            para.get('text', para.get('table', ''))
            for page in result for para in page['content']
        )

    def parse_zip(self, file_path: str) -> List[dict]:
        extract_dir = os.path.join(self.data_root, f"zip_{hash_sha256(file_path)}")
        os.makedirs(extract_dir, exist_ok=True)

        results = []
        for extracted_file in parse_zip(file_path, extract_dir):
            if (ft := get_file_type(extracted_file)) in self.parsers:
                try:
                    results.extend(self.parsers[ft](extracted_file))
                except Exception as e:
                    logger.warning(f"Skip files {extracted_file}: {str(e)}")

        if not results:
            raise ValueError("No parseable content found in the ZIP file")
        return results
